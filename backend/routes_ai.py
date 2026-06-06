"""
AI routes for AcademiaOS Phase 3.
"""

import os
import re
import io
import json
import uuid
import shutil
import logging
from collections import Counter
from datetime import datetime, timezone
from typing import List, Optional

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, Request, Body
from fastapi.responses import FileResponse
from pydantic import BaseModel

import ai_service
from ai_service import (
    resolve_model, generate_text, generate_json, get_or_create_session,
    chat_send, retrieve, chunk_text, _tokens, now_iso,
)

logger = logging.getLogger("academiaos.ai_routes")

UPLOAD_DIR = os.environ.get("UPLOAD_DIR", "/app/uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# File-type constraints
ALLOWED_MIME = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "text/plain": ".txt",
    "text/markdown": ".md",
}
MAX_BYTES = 25 * 1024 * 1024  # 25 MB


def _extract_text(path: str, ext: str) -> str:
    """Best-effort text extraction by file extension. Returns first 60k chars."""
    ext = (ext or "").lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(path)
            parts = []
            for page in reader.pages[:120]:
                try:
                    parts.append(page.extract_text() or "")
                except Exception:
                    continue
            return "\n\n".join(parts)[:60000]
        if ext == ".docx":
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text)[:60000]
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            return "\n".join(chunks)[:60000]
        # txt / md / fallback
        with open(path, "rb") as f:
            return f.read().decode("utf-8", errors="ignore")[:60000]
    except Exception as e:
        logger.exception("Text extraction failed for %s: %s", path, e)
        return ""


def build_router(get_db, get_current_user):
    router = APIRouter(prefix="/api/ai")

    async def _scope(user: dict, institution_id: str):
        if user["role"] != "super_admin" and user.get("institution_id") != institution_id:
            raise HTTPException(403, "Forbidden")

    # ---------------------------------------------------------------------
    # AI Use Cases catalog + config
    # ---------------------------------------------------------------------
    @router.get("/use-cases/{institution_id}")
    async def list_use_cases(institution_id: str, user: dict = Depends(get_current_user)):
        await _scope(user, institution_id)
        db = get_db()
        items = await db.ai_use_cases.find(
            {"institution_id": institution_id}, {"_id": 0}
        ).to_list(50)
        # stable order by code
        items.sort(key=lambda x: x.get("code", ""))
        return items

    class UseCasePatch(BaseModel):
        status: Optional[str] = None
        provider: Optional[str] = None
        model: Optional[str] = None
        human_in_the_loop: Optional[bool] = None
        citations_required: Optional[bool] = None
        risk_score: Optional[str] = None

    @router.patch("/use-cases/{institution_id}/{key}")
    async def patch_use_case(
        institution_id: str, key: str, payload: UseCasePatch = Body(...),
        user: dict = Depends(get_current_user),
    ):
        await _scope(user, institution_id)
        db = get_db()
        updates = {k: v for k, v in payload.model_dump().items() if v is not None}
        if not updates:
            raise HTTPException(400, "No fields to update")
        await db.ai_use_cases.update_one(
            {"institution_id": institution_id, "key": key}, {"$set": updates}
        )
        uc = await db.ai_use_cases.find_one(
            {"institution_id": institution_id, "key": key}, {"_id": 0}
        )
        await db.audit_logs.insert_one({
            "id": str(uuid.uuid4()),
            "institution_id": institution_id,
            "action": "ai.use_case.update",
            "target": key,
            "actor": user["email"],
            "changes": updates,
            "ts": now_iso(),
        })
        return uc

    # ---------------------------------------------------------------------
    # Content Studio — upload + AI generation + SME approval
    # ---------------------------------------------------------------------
    @router.get("/content/sources/{institution_id}")
    async def list_sources(institution_id: str, user: dict = Depends(get_current_user)):
        await _scope(user, institution_id)
        db = get_db()
        items = await db.content_sources.find(
            {"institution_id": institution_id}, {"_id": 0}
        ).sort("uploaded_at", -1).to_list(200)
        return items

    @router.post("/content/upload")
    async def upload_source(
        request: Request,
        institution_id: str = Form(...),
        title: str = Form(...),
        course_id: str = Form(""),
        kind: str = Form("lecture_notes"),
        file: Optional[UploadFile] = File(None),
        text: str = Form(""),
        user: dict = Depends(get_current_user),
    ):
        await _scope(user, institution_id)
        db = get_db()

        source_id = str(uuid.uuid4())
        filename = None
        stored_ext = None
        size_bytes = 0
        mime = None
        body_text = text or ""
        if file is not None:
            # MIME / extension validation
            mime = (file.content_type or "").lower()
            raw_name = file.filename or "upload.bin"
            ext = os.path.splitext(raw_name)[1].lower()
            if mime in ALLOWED_MIME:
                stored_ext = ALLOWED_MIME[mime]
            elif ext in {".pdf", ".docx", ".pptx", ".txt", ".md"}:
                stored_ext = ext
            else:
                raise HTTPException(415, f"Unsupported file type: {mime or ext or 'unknown'}. Allowed: PDF, DOCX, PPTX, TXT, MD.")

            filename = f"{source_id}{stored_ext}"
            target = os.path.join(UPLOAD_DIR, filename)
            # streamed write with size cap
            with open(target, "wb") as out:
                while True:
                    chunk = await file.read(1024 * 64)
                    if not chunk:
                        break
                    size_bytes += len(chunk)
                    if size_bytes > MAX_BYTES:
                        out.close()
                        try:
                            os.remove(target)
                        except OSError:
                            pass
                        raise HTTPException(413, f"File too large (max {MAX_BYTES // (1024*1024)} MB)")
                    out.write(chunk)
            # type-aware extraction
            body_text = _extract_text(target, stored_ext)
            if not body_text and not (text or "").strip():
                logger.warning("No text extracted from %s — proceeding with empty body", filename)

        doc = {
            "id": source_id,
            "institution_id": institution_id,
            "course_id": course_id or None,
            "title": title,
            "kind": kind,
            "filename": filename,
            "original_filename": (file.filename if file else None),
            "mime": mime,
            "size_bytes": size_bytes,
            "text": body_text,
            "uploaded_by": user.get("name"),
            "uploaded_at": now_iso(),
            "approved": False,
        }
        await db.content_sources.insert_one(dict(doc))
        await db.audit_logs.insert_one({
            "id": str(uuid.uuid4()),
            "institution_id": institution_id,
            "action": "content.upload",
            "target": source_id,
            "actor": user["email"],
            "size_bytes": size_bytes,
            "mime": mime,
            "ts": now_iso(),
        })
        doc.pop("_id", None)
        return doc

    @router.get("/content/sources/{source_id}/download")
    async def download_source(source_id: str, user: dict = Depends(get_current_user)):
        db = get_db()
        src = await db.content_sources.find_one({"id": source_id}, {"_id": 0})
        if not src:
            raise HTTPException(404, "Source not found")
        await _scope(user, src["institution_id"])
        if not src.get("filename"):
            raise HTTPException(404, "This source has no uploaded file")
        path = os.path.join(UPLOAD_DIR, src["filename"])
        if not os.path.isfile(path):
            raise HTTPException(404, "Stored file missing on disk")
        return FileResponse(
            path,
            media_type=src.get("mime") or "application/octet-stream",
            filename=src.get("original_filename") or src["filename"],
        )

    @router.post("/content/{source_id}/approve")
    async def approve_source(source_id: str, user: dict = Depends(get_current_user)):
        db = get_db()
        src = await db.content_sources.find_one({"id": source_id}, {"_id": 0})
        if not src:
            raise HTTPException(404, "Source not found")
        await _scope(user, src["institution_id"])
        await db.content_sources.update_one(
            {"id": source_id},
            {"$set": {"approved": True, "approved_by": user["name"], "approved_at": now_iso()}},
        )
        # Index chunks for RAG
        chunks = chunk_text(src.get("text", ""))
        await db.content_chunks.delete_many({"source_id": source_id})
        if chunks:
            await db.content_chunks.insert_many([
                {
                    "id": str(uuid.uuid4()),
                    "source_id": source_id,
                    "institution_id": src["institution_id"],
                    "course_id": src.get("course_id"),
                    "title": src["title"],
                    "ordinal": i,
                    "text": c,
                    "tokens": dict(Counter(_tokens(c))),
                    "approved": True,
                }
                for i, c in enumerate(chunks)
            ])
        await db.audit_logs.insert_one({
            "id": str(uuid.uuid4()),
            "institution_id": src["institution_id"],
            "action": "content.approve",
            "target": source_id,
            "actor": user["email"],
            "ts": now_iso(),
        })
        return {"ok": True, "chunks_indexed": len(chunks)}

    class GenerateRequest(BaseModel):
        institution_id: str
        source_id: str
        kind: str  # 'lesson_plan' | 'flashcards' | 'mcqs' | 'case_guide'
        difficulty: str = "intermediate"
        language: str = "en"
        count: int = 5
        bloom: str = "Apply"

    @router.post("/content/generate")
    async def generate_content(req: GenerateRequest = Body(...), user: dict = Depends(get_current_user)):
        await _scope(user, req.institution_id)
        db = get_db()
        src = await db.content_sources.find_one({"id": req.source_id}, {"_id": 0})
        if not src:
            raise HTTPException(404, "Source not found")
        provider, model = await resolve_model(db, req.institution_id, "content_studio")

        sys = (
            "You are AcademiaOS Content Studio — a senior instructional designer. "
            "Generate teaching content grounded ONLY in the supplied source material. "
            f"Output language: {'Arabic (formal modern)' if req.language=='ar' else 'English'}. "
            f"Difficulty: {req.difficulty}. Bloom's level: {req.bloom}. "
            "Always include a 'source_citations' field listing the exact phrases or sentences "
            "from the source that justify each item."
        )

        schema_hint = {
            "lesson_plan": (
                "Schema: {title, overview, learning_outcomes:[string], "
                "session_plan:[{minute, activity, instructor_note}], "
                "key_concepts:[{term, definition}], assessment_prompt, source_citations:[string]}"
            ),
            "flashcards": (
                f"Schema: {{cards:[{{q, a, hint, source_citation}}] of length {req.count}}}"
            ),
            "mcqs": (
                f"Schema: {{questions:[{{stem, options:[A,B,C,D], correct_index:int, explanation, "
                f"bloom, source_citation}}] of length {req.count}}}"
            ),
            "case_guide": (
                "Schema: {case_title, scenario, discussion_questions:[string], "
                "teaching_note, decision_points:[string], source_citations:[string]}"
            ),
        }.get(req.kind, "Schema: free-form JSON object.")

        user_text = (
            f"### SOURCE TITLE\n{src['title']}\n\n"
            f"### SOURCE TEXT\n{(src.get('text') or '')[:6000]}\n\n"
            f"### TASK\nProduce a {req.kind} based strictly on the source above. {schema_hint}"
        )

        try:
            data = await generate_json(
                system_message=sys, user_text=user_text, provider=provider, model=model
            )
        except Exception as e:
            logger.exception("AI generation failed")
            raise HTTPException(502, f"AI provider error: {e}")

        out_id = str(uuid.uuid4())
        out = {
            "id": out_id,
            "institution_id": req.institution_id,
            "source_id": req.source_id,
            "source_title": src["title"],
            "course_id": src.get("course_id"),
            "kind": req.kind,
            "difficulty": req.difficulty,
            "bloom": req.bloom,
            "language": req.language,
            "model": f"{provider}/{model}",
            "payload": data,
            "status": "pending_review",
            "created_by": user["name"],
            "created_at": now_iso(),
        }
        await db.ai_outputs.insert_one(dict(out))
        await db.audit_logs.insert_one({
            "id": str(uuid.uuid4()),
            "institution_id": req.institution_id,
            "action": "ai.content.generate",
            "target": out_id,
            "actor": user["email"],
            "model": out["model"],
            "ts": now_iso(),
        })
        out.pop("_id", None)
        return out

    @router.get("/content/outputs/{institution_id}")
    async def list_outputs(institution_id: str, status: Optional[str] = None, user: dict = Depends(get_current_user)):
        await _scope(user, institution_id)
        db = get_db()
        q = {"institution_id": institution_id}
        if status:
            q["status"] = status
        items = await db.ai_outputs.find(q, {"_id": 0}).sort("created_at", -1).to_list(200)
        return items

    @router.post("/content/outputs/{output_id}/approve")
    async def approve_output(output_id: str, user: dict = Depends(get_current_user)):
        db = get_db()
        out = await db.ai_outputs.find_one({"id": output_id}, {"_id": 0})
        if not out:
            raise HTTPException(404, "Not found")
        await _scope(user, out["institution_id"])
        await db.ai_outputs.update_one(
            {"id": output_id},
            {"$set": {"status": "approved", "approved_by": user["name"], "approved_at": now_iso()}},
        )
        await db.audit_logs.insert_one({
            "id": str(uuid.uuid4()),
            "institution_id": out["institution_id"],
            "action": "ai.content.approve",
            "target": output_id,
            "actor": user["email"],
            "ts": now_iso(),
        })
        return {"ok": True}

    @router.post("/content/outputs/{output_id}/reject")
    async def reject_output(output_id: str, user: dict = Depends(get_current_user)):
        db = get_db()
        out = await db.ai_outputs.find_one({"id": output_id}, {"_id": 0})
        if not out:
            raise HTTPException(404, "Not found")
        await _scope(user, out["institution_id"])
        await db.ai_outputs.update_one(
            {"id": output_id}, {"$set": {"status": "rejected", "rejected_by": user["name"]}}
        )
        return {"ok": True}

    # ---------------------------------------------------------------------
    # AI Instructor (course-grounded chat)
    # ---------------------------------------------------------------------
    class InstructorMsg(BaseModel):
        institution_id: str
        course_id: Optional[str] = None
        text: str
        language: str = "en"
        persona: Optional[str] = "lecturer"  # lecturer | tutor | coach | examiner
        depth: Optional[str] = "standard"    # concise | standard | deep
        show_reasoning: Optional[bool] = False

    PERSONAS = {
        "lecturer": "You are a structured university lecturer. Open with a one-sentence definition, then explain the core idea with one concrete example, and close with a callback to the source.",
        "tutor": "You are a one-on-one tutor. Ask one clarifying question if intent is ambiguous, otherwise teach with a worked micro-example and 1-2 follow-up prompts.",
        "coach": "You are an executive coach. Be direct, action-oriented and end with a specific next action the learner should take this week.",
        "examiner": "You are an examiner preparing the learner for an assessment. Give the answer, then quiz the learner with one short retrieval question.",
    }
    DEPTHS = {
        "concise": "Keep the answer under 90 words.",
        "standard": "Aim for 140-220 words with 1 example.",
        "deep": "Provide a 280-420 word response with 2 examples and an edge case.",
    }

    @router.post("/instructor/message")
    async def instructor_message(req: InstructorMsg = Body(...), user: dict = Depends(get_current_user)):
        await _scope(user, req.institution_id)
        db = get_db()
        provider, model = await resolve_model(db, req.institution_id, "ai_instructor")

        # Retrieve context
        t0 = datetime.now(timezone.utc)
        hits = await retrieve(
            db, institution_id=req.institution_id, query=req.text, top_k=4, course_id=req.course_id
        )
        # Fallback: if no course-scoped chunks match, retry tenant-wide so
        # obvious questions still surface an approved source instead of
        # showing an empty citations rail.
        if not hits and req.course_id:
            hits = await retrieve(
                db, institution_id=req.institution_id, query=req.text, top_k=4, course_id=None
            )
        context_block = "\n\n".join(
            f"[{i+1}] ({h['title']}) {h['text']}" for i, h in enumerate(hits)
        )

        course_label = ""
        if req.course_id:
            course = await db.courses.find_one({"id": req.course_id}, {"_id": 0})
            if course:
                course_label = f" for the course '{course['title']}'"

        lang = "Arabic (formal modern)" if req.language == "ar" else "English"
        persona_directive = PERSONAS.get(req.persona or "lecturer", PERSONAS["lecturer"])
        depth_directive = DEPTHS.get(req.depth or "standard", DEPTHS["standard"])
        reasoning_directive = (
            "Begin with a <reasoning>…</reasoning> block summarising in ≤40 words which sources you'll use and why. "
            "Then answer normally."
            if req.show_reasoning else ""
        )
        sys = (
            f"You are the AcademiaOS Virtual AI Instructor{course_label}. Respond in {lang}. "
            f"{persona_directive} {depth_directive} {reasoning_directive} "
            "Use ONLY the approved course material between <SOURCE> tags below. "
            "If the answer is not in the source, say so and offer to escalate to faculty. "
            "After your answer, list citations as [1], [2] mapped to the bracketed source numbers."
            "\n\n<SOURCE>\n" + (context_block or "(no sources retrieved)") + "\n</SOURCE>"
        )

        sess = await get_or_create_session(
            db,
            institution_id=req.institution_id,
            user_id=user["id"],
            kind="instructor",
            course_id=req.course_id,
        )
        try:
            citations = [
                {"n": i + 1, "title": h["title"], "score": h["score"], "snippet": h["text"][:200]}
                for i, h in enumerate(hits)
            ]
            assistant = await chat_send(
                db, session=sess, user_text=req.text,
                system_message=sys, provider=provider, model=model,
                citations=citations,
            )
            latency_ms = int((datetime.now(timezone.utc) - t0).total_seconds() * 1000)
            # Extract reasoning block if requested
            reply_text = assistant["text"]
            reasoning = None
            if req.show_reasoning and "<reasoning>" in reply_text:
                try:
                    start = reply_text.index("<reasoning>") + len("<reasoning>")
                    end = reply_text.index("</reasoning>")
                    reasoning = reply_text[start:end].strip()
                    reply_text = (reply_text[:reply_text.index("<reasoning>")] + reply_text[end + len("</reasoning>"):]).strip()
                except ValueError:
                    pass
            return {
                "reply": reply_text,
                "citations": citations,
                "model": assistant["model"],
                "session_id": sess["id"],
                "latency_ms": latency_ms,
                "persona": req.persona,
                "depth": req.depth,
                "reasoning": reasoning,
                "tokens_in": assistant.get("tokens_in"),
                "tokens_out": assistant.get("tokens_out"),
            }
        except Exception as e:
            logger.exception("Instructor chat failed")
            raise HTTPException(502, f"AI error: {e}")

    @router.get("/instructor/suggestions/{institution_id}")
    async def instructor_suggestions(
        institution_id: str,
        course_id: Optional[str] = None,
        language: str = "en",
        user: dict = Depends(get_current_user),
    ):
        """Returns 4 tenant- and course-aware starter prompts.

        Static — fast, deterministic and avoids spending LLM budget on UI hints.
        """
        await _scope(user, institution_id)
        db = get_db()
        institution = await db.institutions.find_one({"id": institution_id}, {"_id": 0}) or {}
        course = None
        if course_id:
            course = await db.courses.find_one({"id": course_id}, {"_id": 0})
        title = (course or {}).get("title") or institution.get("name") or "this course"
        if language == "ar":
            return {
                "items": [
                    f"اشرح الفكرة الأساسية في {title} بأسلوب موجز.",
                    "ما هي أهم المخرجات التعليمية لهذه الوحدة؟",
                    "أعطني سؤال مراجعة لاختبار فهمي.",
                    "ما الفرق بين النهج النظري والتطبيقي في هذا المجال؟",
                ]
            }
        # English defaults — slightly tailored by institution country if available.
        # We retrieve the top approved source titles for this course/tenant so the
        # suggested prompts share vocabulary with the indexed text (better RAG hit-rate).
        src_filter = {"institution_id": institution_id, "approved": True}
        if course_id:
            src_filter["$or"] = [{"course_id": course_id}, {"course_id": None}]
        top_sources = await db.content_sources.find(
            src_filter, {"_id": 0, "title": 1}
        ).sort("uploaded_at", -1).limit(3).to_list(3)
        anchor = (top_sources[0]["title"] if top_sources else title)
        examples = [
            f"Give me a 90-second elevator pitch on {anchor}.",
            f"What are the core concepts in {anchor}? Quiz me on one.",
            "Walk me through the main idea with one concrete example, cited.",
        ]
        country = (institution.get("country") or "").lower()
        if "emirates" in country or "uae" in country:
            examples.append("Frame this concept with a UAE federal-sector example.")
        elif "india" in country:
            examples.append("Apply this to an Indian business context — give one mini case.")
        else:
            examples.append("Show me how this idea applies in industry today, with one example.")
        return {"items": examples}

    @router.get("/instructor/sessions/{institution_id}")
    async def list_instructor_sessions(institution_id: str, user: dict = Depends(get_current_user)):
        await _scope(user, institution_id)
        db = get_db()
        items = await db.ai_sessions.find(
            {"institution_id": institution_id, "kind": "instructor"}, {"_id": 0}
        ).sort("created_at", -1).limit(40).to_list(40)
        return items

    # ---------------------------------------------------------------------
    # AI Advisor — skill gap + learning path
    # ---------------------------------------------------------------------
    @router.get("/advisor/framework/{institution_id}")
    async def get_framework(institution_id: str, user: dict = Depends(get_current_user)):
        await _scope(user, institution_id)
        db = get_db()
        fw = await db.skill_frameworks.find_one({"institution_id": institution_id}, {"_id": 0})
        return fw or {"institution_id": institution_id, "target_roles": []}

    @router.get("/advisor/profile/{institution_id}/{user_id}")
    async def get_profile(institution_id: str, user_id: str, user: dict = Depends(get_current_user)):
        await _scope(user, institution_id)
        db = get_db()
        p = await db.learner_profiles.find_one(
            {"institution_id": institution_id, "user_id": user_id}, {"_id": 0}
        )
        if not p:
            raise HTTPException(404, "Profile not found")
        return p

    class AdvisorRequest(BaseModel):
        institution_id: str
        user_id: str
        language: str = "en"

    @router.post("/advisor/analyse")
    async def advisor_analyse(req: AdvisorRequest = Body(...), user: dict = Depends(get_current_user)):
        await _scope(user, req.institution_id)
        db = get_db()
        profile = await db.learner_profiles.find_one(
            {"institution_id": req.institution_id, "user_id": req.user_id}, {"_id": 0}
        )
        if not profile:
            raise HTTPException(404, "Profile not found")
        fw = await db.skill_frameworks.find_one({"institution_id": req.institution_id}, {"_id": 0}) or {}
        target_role = next(
            (r for r in fw.get("target_roles", []) if r["key"] == profile.get("target_role")),
            None,
        )
        if not target_role:
            raise HTTPException(404, "Target role not configured")

        # Compute deterministic skill gap
        current = {s["name"]: s["level"] for s in profile.get("skills", [])}
        gaps = []
        for s in target_role["skills"]:
            cur = current.get(s["name"], 0)
            gaps.append({
                "skill": s["name"],
                "current": cur,
                "target": s["level"],
                "gap": max(0, s["level"] - cur),
            })

        provider, model = await resolve_model(db, req.institution_id, "ai_advisor")
        learner = await db.users.find_one({"id": req.user_id}, {"_id": 0, "password_hash": 0}) or {}
        # available courses
        courses = await db.courses.find({"institution_id": req.institution_id}, {"_id": 0}).to_list(80)
        course_catalog = "\n".join(f"- {c['title']} ({c['code']})" for c in courses[:30])

        lang = "Arabic (formal modern)" if req.language == "ar" else "English"
        sys = (
            f"You are the AcademiaOS Educational Advisor. Respond in {lang}. "
            "Given a learner profile, target role and a course catalog, produce a JSON "
            "recommendation. Be concrete, encouraging and time-bounded."
        )
        user_text = (
            f"LEARNER: {learner.get('name')} — {learner.get('title')}\n"
            f"TARGET ROLE: {target_role['name']}\n"
            f"SKILL GAPS (skill, current, target, gap):\n"
            + "\n".join(f"- {g['skill']}: {g['current']} → {g['target']} (gap {g['gap']})" for g in gaps)
            + f"\n\nCOURSE CATALOG (institution offerings):\n{course_catalog}\n\n"
            "Return JSON with schema: {summary, top_priorities:[{skill, why}], "
            "recommended_path:[{order, course_title, duration_weeks, why}], "
            "career_pathway:[{stage, role, horizon}], proactive_alerts:[string]}"
        )
        try:
            payload = await generate_json(
                system_message=sys, user_text=user_text, provider=provider, model=model
            )
        except Exception as e:
            logger.exception("Advisor analyse failed")
            raise HTTPException(502, f"AI error: {e}")

        out_id = str(uuid.uuid4())
        record = {
            "id": out_id,
            "institution_id": req.institution_id,
            "user_id": req.user_id,
            "target_role": target_role,
            "gaps": gaps,
            "payload": payload,
            "model": f"{provider}/{model}",
            "created_at": now_iso(),
        }
        await db.advisor_reports.insert_one(dict(record))
        record.pop("_id", None)
        return record

    # ---------------------------------------------------------------------
    # AI Student Assistant — VEDA chat entry point.
    # Multi-role (student / faculty / admin / parent), multi-lingual
    # (en / hi / te / ar with code-switching awareness), and RAG-grounded
    # against approved Content Studio sources.
    # ---------------------------------------------------------------------
    class AssistantMsg(BaseModel):
        institution_id: str
        text: str
        language: str = "en"  # en | hi | te | ar
        # When omitted, the assistant infers persona from `user["role"]`.
        role_override: Optional[str] = None

    FAQ = (
        "AcademiaOS academic services FAQ:\n"
        "- Enrolment status: visible under Student Dashboard → My Programmes.\n"
        "- Timetable: under Student Dashboard → Schedule, exported as iCal.\n"
        "- Attendance policy: minimum 75% per term, exceptions need programme manager approval.\n"
        "- Assignments and deadlines: under each course module.\n"
        "- Exam dates: published 4 weeks ahead in the academic calendar.\n"
        "- Certificates: issued automatically after final assessment is approved by the registrar.\n"
        "- Graduation requirements: all required modules approved + capstone + credit threshold.\n"
        "- Student support: open a ticket and the assistant escalates to programme office within 24h SLA."
    )

    # Role-specific persona blocks for VEDA
    _PERSONA = {
        "student":
            "You are responding to a STUDENT. Prioritise: timetable, attendance, fees, exams, "
            "certificates, placement readiness, mental-health resources. Be encouraging and concise.",
        "faculty":
            "You are responding to a FACULTY member. Prioritise: course allocation, attendance "
            "submission, research publications, OBE attainment, IQAC processes. Be precise and "
            "reference policy / NAAC / NBA where relevant.",
        "admin":
            "You are responding to an ADMINISTRATOR. Prioritise: cohort analytics, finance, "
            "compliance (NIRF/NAAC/AQAR), staff workload, accreditation timelines. "
            "Provide actionable summaries with named owners.",
        "parent":
            "You are responding to a PARENT/GUARDIAN. Prioritise: attendance %, fee schedule, "
            "exam dates, hostel contact, ward's wellbeing. Use plain language, avoid jargon, "
            "and clearly state where the parent should escalate (programme office vs hostel warden).",
    }
    _ROLE_TO_PERSONA = {
        "student": "student", "guardian": "parent", "parent": "parent",
        "faculty": "faculty", "instructor": "faculty",
        "institution_admin": "admin", "super_admin": "admin",
        "registrar": "admin", "compliance_officer": "admin",
        "ai_governance_admin": "admin", "career_services": "admin",
        "programme_manager": "admin", "training_manager": "admin",
        "hostel_warden": "admin", "counselor": "admin",
    }
    _LANG_LABEL = {
        "en": "English",
        "hi": "Hindi (Devanagari script)",
        "te": "Telugu (Telugu script)",
        "ar": "Arabic (formal modern, RTL)",
    }

    @router.post("/assistant/message")
    async def assistant_message(req: AssistantMsg = Body(...), user: dict = Depends(get_current_user)):
        await _scope(user, req.institution_id)
        db = get_db()
        provider, model = await resolve_model(db, req.institution_id, "student_assistant")
        inst = await db.institutions.find_one({"id": req.institution_id}, {"_id": 0}) or {}

        # 1) Resolve persona from explicit override or user.role
        persona_key = req.role_override or _ROLE_TO_PERSONA.get(user.get("role"), "student")
        persona_key = persona_key if persona_key in _PERSONA else "student"
        persona_block = _PERSONA[persona_key]

        # 2) Resolve language with code-switch instruction
        lang_label = _LANG_LABEL.get(req.language, "English")
        code_switch = (
            "Mirror the user's language. If they code-switch (mix English with Hindi/Telugu), "
            "respond in the same code-switched style they used."
        ) if req.language != "en" else ""

        # 3) RAG retrieval over approved Content Studio sources
        passages = []
        try:
            passages = await retrieve(
                db, institution_id=req.institution_id, query=req.text, top_k=4,
            )
        except Exception:
            passages = []
        rag_block = (
            "\n\n<KNOWLEDGE_BASE>\n" + "\n\n".join(
                f"[Doc {i+1} — {p.get('source_title') or 'untitled'}]\n{p.get('text','')[:480]}"
                for i, p in enumerate(passages)
            ) + "\n</KNOWLEDGE_BASE>"
        ) if passages else ""

        sys = (
            f"You are VEDA, the AcademiaOS assistant for {inst.get('name','this institution')}. "
            f"Respond in {lang_label}. {code_switch}\n\n"
            f"PERSONA: {persona_block}\n\n"
            "GROUNDING RULES:\n"
            "1. If a KNOWLEDGE_BASE block is present, use it as the AUTHORITATIVE source and cite "
            "doc numbers like [Doc 1] inline.\n"
            "2. Otherwise fall back to the static FAQ.\n"
            "3. If neither covers the question, say so politely and offer to open a support ticket.\n"
            "4. Stay in role-scope: a student persona must NOT discuss admin/finance reports."
            "\n\n<FAQ>\n" + FAQ + "</FAQ>"
            + rag_block
        )

        sess = await get_or_create_session(
            db, institution_id=req.institution_id, user_id=user["id"], kind="assistant"
        )
        try:
            assistant = await chat_send(
                db, session=sess, user_text=req.text,
                system_message=sys, provider=provider, model=model,
                citations=[{"source_id": p.get("source_id"), "title": p.get("source_title"),
                            "score": p.get("score")} for p in passages],
                max_history=20,
            )
            return {
                "reply": assistant["text"], "model": assistant["model"],
                "session_id": sess["id"],
                "persona": persona_key, "language": req.language,
                "grounding": "rag" if passages else "faq",
                "citations": assistant.get("citations", []),
            }
        except Exception as e:
            logger.exception("Assistant chat failed")
            raise HTTPException(502, f"AI error: {e}")

    @router.post("/sessions/{session_id}/reset")
    async def reset_session(session_id: str, user: dict = Depends(get_current_user)):
        db = get_db()
        await db.ai_sessions.update_one({"id": session_id}, {"$set": {"open": False}})
        return {"ok": True}

    # ---------------------------------------------------------------------
    # Audit feed (used by Dashboard / Compliance later)
    # ---------------------------------------------------------------------
    @router.get("/audit/{institution_id}")
    async def audit(institution_id: str, limit: int = 50, user: dict = Depends(get_current_user)):
        await _scope(user, institution_id)
        db = get_db()
        items = await db.audit_logs.find(
            {"institution_id": institution_id}, {"_id": 0}
        ).sort("ts", -1).limit(limit).to_list(limit)
        return items

    return router
